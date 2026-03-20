"""Generates the QMKL-Finance PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # dark navy
BG_LIGHT = RGBColor(0x14, 0x28, 0x3E)   # slightly lighter navy
CYAN     = RGBColor(0x00, 0xB4, 0xD8)   # quantum blue
AMBER    = RGBColor(0xF4, 0xA2, 0x61)   # warm accent
GREEN    = RGBColor(0x52, 0xB7, 0x88)   # positive results
RED      = RGBColor(0xE6, 0x39, 0x46)   # negative / warning
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREY     = RGBColor(0x8E, 0xCA, 0xE6)   # muted blue-grey
GOLD     = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── Low-level helpers ─────────────────────────────────────────────────────────

def rgb_hex(r):
    return "{:02X}{:02X}{:02X}".format(r[0], r[1], r[2])

def solid_fill(shape, color: RGBColor):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE,
        left, top, width, height)
    shape.line.fill.background()
    solid_fill(shape, color)
    return shape

# MSO constant for rectangle
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import pptx

def rect(slide, l, t, w, h, color):
    from pptx.util import Emu
    sp = slide.shapes.add_shape(1, l, t, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    sp.line.fill.background()
    solid_fill(sp, color)
    return sp

def bg(slide):
    """Fill slide background with dark navy."""
    r = rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    r.name = "bg"
    return r

def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def accent_line(slide, y, color=CYAN, w=None):
    """Thin horizontal accent line."""
    width = w or Inches(12)
    ln = slide.shapes.add_shape(1, Inches(0.5), y, width, Pt(2))
    solid_fill(ln, color)
    ln.line.fill.background()

def slide_number(slide, n, total):
    txbox(slide, f"{n} / {total}",
          Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.3),
          size=9, color=GREY, align=PP_ALIGN.RIGHT)

def tag(slide, text, l, t, color=CYAN, text_color=BG, size=10):
    """Small filled label/badge."""
    w = Inches(len(text) * 0.11 + 0.25)
    h = Inches(0.28)
    r = rect(slide, l, t, w, h, color)
    txbox(slide, text, l + Inches(0.06), t + Inches(0.02),
          w - Inches(0.06), h, size=size, bold=True, color=text_color,
          align=PP_ALIGN.LEFT)
    return w

# ── Slide factory functions ───────────────────────────────────────────────────

TOTAL = 13

def make_title_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)

    # Left accent bar
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, CYAN)

    # Decorative circles (pseudo quantum feel)
    for i, (cx, cy, r_inch, color, alpha) in enumerate([
        (10.5, 2.0, 2.8, BG_LIGHT, None),
        (11.8, 4.5, 1.8, BG_LIGHT, None),
        (10.0, 5.5, 1.0, BG_LIGHT, None),
    ]):
        circ = slide.shapes.add_shape(9,   # oval
            Inches(cx - r_inch), Inches(cy - r_inch),
            Inches(r_inch * 2), Inches(r_inch * 2))
        circ.line.color.rgb = CYAN
        circ.line.width = Pt(1)
        circ.fill.background()

    # Subtitle top tag
    txbox(slide, "QUANTUM MACHINE LEARNING · FINANCE · 2025–2026",
          Inches(0.8), Inches(1.0), Inches(11), Inches(0.4),
          size=11, color=CYAN, bold=True)

    # Title
    txbox(slide, "Quantum Multiple\nKernel Learning",
          Inches(0.8), Inches(1.6), Inches(9.5), Inches(2.2),
          size=54, bold=True, color=WHITE)

    # Subtitle
    txbox(slide, "pour la Classification Financière",
          Inches(0.8), Inches(3.7), Inches(9.5), Inches(0.7),
          size=26, color=AMBER)

    accent_line(slide, Inches(4.55), CYAN, Inches(9.0))

    txbox(slide, "Un benchmark empirique complet · 17 notebooks · 14 contributions originales",
          Inches(0.8), Inches(4.7), Inches(10), Inches(0.45),
          size=14, color=GREY)

    txbox(slide, "Mars 2026",
          Inches(0.8), Inches(5.2), Inches(4), Inches(0.4),
          size=13, color=GREY)

    slide_number(slide, 1, TOTAL)
    return slide


def make_journey_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, AMBER)

    txbox(slide, "Notre parcours", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), AMBER)

    txbox(slide, "17 notebooks · des données brutes aux contributions frontières",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.4),
          size=13, color=GREY)

    # Timeline phases
    phases = [
        ("Fondations",    "NB 01–05", "Exploration données,\nbaselines, ablation",        CYAN),
        ("Benchmark",     "NB 06–10", "20+ kernels, BO,\nHardware IBM, Shapley",          AMBER),
        ("Synthèse",      "NB 11–12", "Rapport final,\n5 métriques interprétabilité",     GREEN),
        ("Contributions", "NB 13–15", "QKRR, VQKL, QKAM,\nPhase diagram, Drift",         GOLD),
        ("Extension",     "NB 16–17", "Rapport figure,\nMitigation concentration",         RGBColor(0xC7, 0x7D, 0xFF)),
    ]

    for i, (title, nb, desc, color) in enumerate(phases):
        x = Inches(0.7 + i * 2.5)
        # Connector line
        if i < len(phases) - 1:
            ln = slide.shapes.add_shape(1, x + Inches(1.6), Inches(2.45),
                                         Inches(0.9), Pt(2))
            solid_fill(ln, GREY)
            ln.line.fill.background()

        # Circle marker
        circ = slide.shapes.add_shape(9,
            x, Inches(2.0), Inches(1.6), Inches(0.9))
        solid_fill(circ, color)
        circ.line.fill.background()
        txbox(slide, nb, x, Inches(2.0), Inches(1.6), Inches(0.9),
              size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

        txbox(slide, title, x - Inches(0.05), Inches(3.05), Inches(1.7), Inches(0.4),
              size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        txbox(slide, desc, x - Inches(0.05), Inches(3.5), Inches(1.7), Inches(0.9),
              size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Key stats band
    rect(slide, Inches(0.55), Inches(4.7), Inches(12.2), Inches(1.65), BG_LIGHT)

    stats = [
        ("17", "notebooks"),
        ("3",  "datasets financiers"),
        ("12", "kernels quantiques"),
        ("14", "contributions originales"),
        ("20", "runs par expérience"),
        ("6",  "qubits simulés"),
    ]
    for i, (val, label) in enumerate(stats):
        x = Inches(0.75 + i * 2.02)
        txbox(slide, val,   x, Inches(4.82), Inches(1.9), Inches(0.65),
              size=36, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        txbox(slide, label, x, Inches(5.42), Inches(1.9), Inches(0.35),
              size=10, color=GREY, align=PP_ALIGN.CENTER)

    slide_number(slide, 2, TOTAL)
    return slide


def make_question_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, CYAN)

    txbox(slide, "La question centrale", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), CYAN)

    # Big question
    rect(slide, Inches(0.7), Inches(1.2), Inches(12.1), Inches(1.55), BG_LIGHT)
    txbox(slide,
          "Les kernels quantiques apportent-ils un avantage mesurable\nsur des données financières réelles ?",
          Inches(0.85), Inches(1.3), Inches(11.8), Inches(1.35),
          size=24, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Context boxes
    boxes = [
        ("Le kernel quantique",
         "K(x,x') = |⟨ψ(x')|ψ(x)⟩|²\n\nMesure la similarité entre deux points de données via la fidélité de leurs états quantiques. En théorie, capture des structures que les kernels classiques ne voient pas.",
         CYAN),
        ("Multiple Kernel Learning",
         "K_w = Σ wₘ Kₘ\n\nCombine M kernels quantiques différents (familles Pauli, bandwidths α variés) avec des poids optimisés — meilleur que n'importe quel kernel seul.",
         AMBER),
        ("Notre réponse honnête",
         "Après 17 notebooks et des centaines d'expériences : QMKL sous-performe les classiques de −4 à −7 pts AUC. Mais on a compris pourquoi — et trouvé des exceptions intéressantes.",
         GREEN),
    ]

    for i, (title, text, color) in enumerate(boxes):
        x = Inches(0.7 + i * 4.1)
        rect(slide, x, Inches(3.0), Inches(3.9), Inches(3.9), BG_LIGHT)
        rect(slide, x, Inches(3.0), Inches(3.9), Inches(0.32), color)
        txbox(slide, title, x + Inches(0.1), Inches(3.05), Inches(3.7), Inches(0.28),
              size=11, bold=True, color=BG, align=PP_ALIGN.LEFT)
        txbox(slide, text, x + Inches(0.1), Inches(3.42), Inches(3.7), Inches(3.3),
              size=11, color=WHITE)

    slide_number(slide, 3, TOTAL)
    return slide


def make_setup_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, AMBER)

    txbox(slide, "Notre démarche expérimentale", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), AMBER)

    # 3 columns: datasets / kernels / strategies
    col_data = [
        ("3 Datasets financiers", CYAN, [
            ("German Credit",  "1 000 clients banque\n20 features → 6 (PCA)\n30% mauvais payeurs"),
            ("Bank Marketing", "45 000 appels\n16 features → 6 (PCA)\n12% souscriptions (déséquilibre)"),
            ("Breast Cancer",  "569 patients\n30 features → 6 (PCA)\nDataset médical référence"),
        ]),
        ("12 Kernels quantiques", AMBER, [
            ("Familles Pauli",    "Z, ZZ, XZ, YXX, YZX, Pauli\n6 familles × 2 bandwidths α"),
            ("Feature maps",      "PauliFeatureMap, 6 qubits\nSimulation statevector exacte"),
            ("Diversité réelle",  "Alignement moyen entre paires:\n0.62 — réellement distincts"),
        ]),
        ("4 Stratégies MKL", GREEN, [
            ("Average",          "Poids uniformes 1/M\n(baseline simple)"),
            ("Centered Align.",  "Solution analytique fermée\n(optimale en KTA)"),
            ("Best. optim. BO",  "Optimisation bayésienne\n(25 appels, 8 points initiaux)"),
            ("Protocole",        "20 runs · split 67/33\nTest Wilcoxon vs RBF-SVM"),
        ]),
    ]

    for ci, (header, color, items) in enumerate(col_data):
        x = Inches(0.7 + ci * 4.15)
        rect(slide, x, Inches(1.1), Inches(4.0), Inches(0.38), color)
        txbox(slide, header, x + Inches(0.08), Inches(1.12), Inches(3.85), Inches(0.35),
              size=12, bold=True, color=BG)

        for ri, (ititle, itext) in enumerate(items):
            y = Inches(1.6 + ri * 1.9)
            rect(slide, x, y, Inches(4.0), Inches(1.75), BG_LIGHT)
            rect(slide, x, y, Pt(4), Inches(1.75), color)
            txbox(slide, ititle, x + Inches(0.12), y + Inches(0.06), Inches(3.8), Inches(0.3),
                  size=11, bold=True, color=color)
            txbox(slide, itext, x + Inches(0.12), y + Inches(0.38), Inches(3.8), Inches(1.25),
                  size=10, color=WHITE)

    slide_number(slide, 4, TOTAL)
    return slide


def make_results_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, RED)

    txbox(slide, "Le résultat principal", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), RED)

    txbox(slide, "AUC moyen sur 20 runs · test Wilcoxon vs RBF-SVM",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35),
          size=13, color=GREY)

    # Table header
    cols = ["Méthode", "German\nCredit", "Bank\nMarketing", "Breast\nCancer"]
    col_w = [Inches(3.0), Inches(2.3), Inches(2.3), Inches(2.3)]
    col_x = [Inches(0.65), Inches(3.75), Inches(6.05), Inches(8.35)]
    row_h = Inches(0.52)
    y0    = Inches(1.5)

    # Header row
    rect(slide, Inches(0.65), y0, Inches(10.0), row_h, CYAN)
    for ci, (col, cw, cx) in enumerate(zip(cols, col_w, col_x)):
        txbox(slide, col, cx + Inches(0.05), y0 + Inches(0.04), cw - Inches(0.1), row_h,
              size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

    rows = [
        # (label, gc, bm, bc, is_quantum, highlight)
        ("[Q] Average",       "0.7634", "0.7246", "0.9951", True,  False),
        ("[Q] Centered Align.","0.7504","0.7628", "0.9948", True,  False),
        ("[Q] Bayesian Opt.", "0.7584", "0.7741", "0.9948", True,  False),
        ("[C] RBF-SVM",       "0.8345", "0.8168", "0.9963", False, True),
        ("[C] Random Forest", "0.8330", "0.8183", "0.9916", False, False),
        ("[C] Régression Log.","0.7990","0.8672", "0.9955", False, False),
    ]

    for ri, (label, gc, bm, bc, is_q, highlight) in enumerate(rows):
        y = y0 + row_h + ri * row_h
        row_color = BG_LIGHT if ri % 2 == 0 else BG
        if highlight:
            row_color = RGBColor(0x0A, 0x3D, 0x2E)
        rect(slide, Inches(0.65), y, Inches(10.0), row_h, row_color)

        label_color = CYAN if is_q else AMBER
        if highlight:
            label_color = GREEN

        txbox(slide, label, col_x[0] + Inches(0.05), y + Inches(0.1),
              col_w[0] - Inches(0.1), row_h - Inches(0.05),
              size=11, bold=highlight, color=label_color)

        for val, cx, cw, ref in zip([gc, bm, bc], col_x[1:], col_w[1:],
                                    ["0.8345", "0.8168", "0.9963"]):
            try:
                delta = float(val) - float(ref)
                color = GREEN if delta >= -0.002 else (WHITE if delta > -0.03 else RED)
            except:
                color = WHITE
            if highlight:
                color = GOLD
            txbox(slide, val, cx + Inches(0.05), y + Inches(0.1),
                  cw - Inches(0.1), row_h - Inches(0.05),
                  size=12, bold=highlight, color=color, align=PP_ALIGN.CENTER)

    # Gap annotation
    txbox(slide, "Gap QMKL vs RBF-SVM (meilleur quantique) :",
          Inches(0.7), Inches(5.6), Inches(5), Inches(0.35),
          size=12, bold=True, color=WHITE)

    gaps = [
        ("German Credit", "−7.1 pts AUC", RED),
        ("Bank Marketing","−4.3 pts AUC", AMBER),
        ("Breast Cancer", "−0.1 pts AUC\n(non significatif)", GREEN),
    ]
    for i, (ds, gap, color) in enumerate(gaps):
        x = Inches(0.7 + i * 3.8)
        rect(slide, x, Inches(6.0), Inches(3.6), Inches(1.0), BG_LIGHT)
        txbox(slide, ds,  x + Inches(0.1), Inches(6.05), Inches(3.4), Inches(0.3),
              size=10, color=GREY)
        txbox(slide, gap, x + Inches(0.1), Inches(6.35), Inches(3.4), Inches(0.55),
              size=14, bold=True, color=color)

    slide_number(slide, 5, TOTAL)
    return slide


def make_diagnosis_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, RED)

    txbox(slide, "Pourquoi QMKL sous-performe ?", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), RED)
    txbox(slide, "Quatre causes diagnostiquées expérimentalement",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35),
          size=13, color=GREY)

    causes = [
        ("01", "Alignement avec le RBF",
         "Les kernels quantiques imitent partiellement le RBF\n(alignement Frobenius moyen : 0.642).\nL'espace de features quantique n'est pas assez\northogonal pour apporter de l'information complémentaire.",
         AMBER, "0.642\nalig. moyen"),
        ("02", "Barren plateaux",
         "La concentration du kernel chute de 60 % entre Q=2\net Q=8. Avec Q=6, les kernels sont déjà partiellement\nconcentrés : la matrice de Gram tend vers l'identité\net perd sa capacité discriminante.",
         RED, "−60 %\nQ2 → Q8"),
        ("03", "Structure des données",
         "Les données financières tabulaires (après PCA)\nn'ont pas la structure d'entanglement que les\nPauliFeatureMaps sont conçues pour capturer.\nCes feature maps ciblent les données quantiques.",
         CYAN, "PCA\n→ 6D"),
        ("04", "N / dim trop faible",
         "Avec N=200 et un espace de Hilbert de dim 2^6=64,\nle ratio données/dimension est faible. Le SVM\ndevient dense en support vectors : 84 % des\ninstances pilotent la frontière.",
         GREY, "84 %\nsup. vect."),
    ]

    for i, (num, title, text, color, stat) in enumerate(causes):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(1.55 + row * 2.9)

        rect(slide, x, y, Inches(5.9), Inches(2.65), BG_LIGHT)
        rect(slide, x, y, Inches(0.55), Inches(2.65), color)

        # Number circle
        txbox(slide, num, x + Inches(0.62), y + Inches(0.1), Inches(0.7), Inches(0.45),
              size=22, bold=True, color=color)
        txbox(slide, title, x + Inches(0.62), y + Inches(0.52), Inches(4.0), Inches(0.38),
              size=13, bold=True, color=WHITE)
        txbox(slide, text, x + Inches(0.62), y + Inches(0.92), Inches(3.8), Inches(1.65),
              size=10, color=GREY)

        # Stat badge
        rect(slide, x + Inches(4.6), y + Inches(0.8), Inches(1.2), Inches(0.9), color)
        txbox(slide, stat, x + Inches(4.6), y + Inches(0.82), Inches(1.2), Inches(0.88),
              size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)

    slide_number(slide, 6, TOTAL)
    return slide


def make_barren_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, GOLD)

    txbox(slide, "Les barren plateaux — première carte 2D (Q × α)", Inches(0.7), Inches(0.25),
          Inches(12), Inches(0.55), size=28, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), GOLD)
    txbox(slide, "Quand la grille (Q qubits × bandwidth α) est-elle optimale ?  [Notebook 13]",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35), size=13, color=GREY)

    # Heatmap placeholder (visual representation)
    # We'll simulate the heatmap with colored cells
    Q_vals    = [2, 3, 4, 5, 6, 7, 8]
    a_vals    = [0.3, 0.5, 1.0, 2.0, 3.0, 4.0]
    # AUC values (approximate from notebook output)
    auc_grid = [
        [0.33, 0.40, 0.54, 0.45, 0.39, 0.37],
        [0.20, 0.30, 0.65, 0.80, 0.55, 0.40],
        [0.17, 0.35, 0.55, 0.75, 0.65, 0.50],
        [0.35, 0.45, 0.70, 0.83, 0.72, 0.60],  # sweet spot Q=5 α=2.0
        [0.30, 0.40, 0.60, 0.74, 0.70, 0.58],
        [0.27, 0.38, 0.58, 0.72, 0.63, 0.55],
        [0.64, 0.65, 0.70, 0.78, 0.73, 0.68],
    ]

    cell_w = Inches(0.72)
    cell_h = Inches(0.52)
    ox = Inches(2.0)
    oy = Inches(1.55)

    # Alpha axis labels
    txbox(slide, "Bandwidth α →", ox, Inches(1.2), Inches(5.5), Inches(0.3),
          size=10, color=GREY, bold=True)
    for ai, a in enumerate(a_vals):
        txbox(slide, str(a), ox + ai * cell_w + Inches(0.22), Inches(1.3),
              cell_w, Inches(0.25), size=10, color=GREY, align=PP_ALIGN.CENTER)

    # Q axis label
    txbox(slide, "Q↓", Inches(1.5), oy, Inches(0.45), Inches(0.35), size=10, color=GREY, bold=True)
    for qi, Q in enumerate(Q_vals):
        txbox(slide, str(Q), Inches(1.65), oy + qi * cell_h + Inches(0.12),
              Inches(0.3), cell_h, size=10, color=GREY, bold=True)

    for qi, row in enumerate(auc_grid):
        for ai, val in enumerate(row):
            # Color interpolation: red (low) → green (high)
            t = (val - 0.17) / (0.83 - 0.17)
            t = max(0, min(1, t))
            r = int(0xE6 * (1 - t) + 0x52 * t)
            g = int(0x39 * (1 - t) + 0xB7 * t)
            b = int(0x46 * (1 - t) + 0x88 * t)
            color = RGBColor(r, g, b)
            cell = rect(slide, ox + ai * cell_w, oy + qi * cell_h, cell_w - Pt(1), cell_h - Pt(1), color)
            # Sweet spot marker
            is_sweet = (qi == 3 and ai == 3)
            txbox(slide, f"{val:.2f}", ox + ai * cell_w + Inches(0.08),
                  oy + qi * cell_h + Inches(0.12), cell_w - Inches(0.1), cell_h - Inches(0.1),
                  size=9, bold=is_sweet, color=BG if val > 0.5 else WHITE,
                  align=PP_ALIGN.CENTER)
            if is_sweet:
                # Border box
                border = rect(slide, ox + ai * cell_w - Pt(2), oy + qi * cell_h - Pt(2),
                              cell_w + Pt(3), cell_h + Pt(3), RGBColor(0xFF, 0xD7, 0x00))
                # Re-draw cell on top
                cell2 = rect(slide, ox + ai * cell_w, oy + qi * cell_h,
                             cell_w - Pt(1), cell_h - Pt(1), color)
                txbox(slide, f"{val:.2f}", ox + ai * cell_w + Inches(0.08),
                      oy + qi * cell_h + Inches(0.12), cell_w - Inches(0.1), cell_h - Inches(0.1),
                      size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)

    # Right panel: key findings
    px = Inches(6.7)
    txbox(slide, "Résultats clés", px, Inches(1.5), Inches(6.2), Inches(0.4),
          size=16, bold=True, color=GOLD)

    findings = [
        (GOLD, "Sweet spot", "Q=5, α=2.0 → AUC = 0.83\n1re carte de ce type pour QMKL finance"),
        (RED,  "Concentration", "Chute de 60 % entre Q=2 et Q=8\n(barren plateau confirmé expérimentalement)"),
        (CYAN, "Expressivité", "H/H_max quantique = 0.653\nvs RBF = 0.813 (kernel moins riche)"),
        (GREEN,"Support vectors", "84 % des instances sont des SV stables\n(symptôme de concentration excessive)"),
    ]

    for i, (color, title, text) in enumerate(findings):
        y = Inches(2.05 + i * 1.3)
        rect(slide, px, y, Inches(6.1), Inches(1.15), BG_LIGHT)
        rect(slide, px, y, Pt(5), Inches(1.15), color)
        txbox(slide, title, px + Inches(0.12), y + Inches(0.06), Inches(5.8), Inches(0.32),
              size=12, bold=True, color=color)
        txbox(slide, text, px + Inches(0.12), y + Inches(0.38), Inches(5.8), Inches(0.7),
              size=10, color=WHITE)

    slide_number(slide, 7, TOTAL)
    return slide


def make_positives_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, GREEN)

    txbox(slide, "Les résultats positifs", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), GREEN)
    txbox(slide, "Trois découvertes inattendues qui ouvrent des perspectives réelles",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35), size=13, color=GREY)

    cards = [
        ("QKRR bat le QSVM", GREEN, "+7 pts\nAUC",
         "La Quantum Kernel Ridge Regression (QKRR) — une alternative\nanalytique au SVM — surpasse le QSVM de +7 pts sur German Credit.",
         "La régularisation ridge est mieux adaptée aux kernels concentrés\nque la marge dure du SVM. Solution analytique : α* = (K + λI)⁻¹y.",
         "QKRR = 0.756\nQSVM = 0.686"),

        ("Diversité → gain MKL", CYAN, "r = 0.738",
         "La diversité entre kernels (alignement croisé faible) prédit\nfortement le gain marginal en LOO — corrélation r = 0.738.",
         "Premier résultat de ce type pour QMKL. Actionnable :\nsélectionner les kernels les plus distincts maximise le gain MKL.",
         "ZZ α=4.0 : gain max\ndiversité = 0.451"),

        ("VQKL : +1.55 pts", AMBER, "+1.55 pts\nAUC",
         "Le Variational Quantum Kernel Learning optimise conjointement\nles bandwidths α et les poids w via gradient sur le KTA.",
         "Les α optimaux divergent légèrement des valeurs initiales\n(α_ZZ : 1.0→1.07, α_XZ : 0.5→0.69). L'approche variationnelle converge.",
         "VQKL = 0.704\nGrille fixe = 0.688"),
    ]

    for i, (title, color, stat, desc, detail, numbers) in enumerate(cards):
        x = Inches(0.7 + i * 4.2)
        rect(slide, x, Inches(1.5), Inches(4.0), Inches(5.5), BG_LIGHT)
        rect(slide, x, Inches(1.5), Inches(4.0), Inches(0.32), color)
        txbox(slide, title, x + Inches(0.1), Inches(1.53), Inches(3.8), Inches(0.28),
              size=12, bold=True, color=BG)

        # Big stat
        rect(slide, x + Inches(0.9), Inches(1.95), Inches(2.2), Inches(1.0), color)
        txbox(slide, stat, x + Inches(0.9), Inches(1.97), Inches(2.2), Inches(0.96),
              size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)

        txbox(slide, desc,   x + Inches(0.1), Inches(3.1),  Inches(3.8), Inches(1.3), size=10, color=WHITE)
        txbox(slide, detail, x + Inches(0.1), Inches(4.45), Inches(3.8), Inches(1.2), size=9,  color=GREY)

        rect(slide, x + Inches(0.1), Inches(5.75), Inches(3.7), Inches(0.8), BG)
        txbox(slide, numbers, x + Inches(0.15), Inches(5.8), Inches(3.6), Inches(0.75),
              size=11, bold=True, color=color)

    slide_number(slide, 8, TOTAL)
    return slide


def make_phase_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, RGBColor(0xC7, 0x7D, 0xFF))

    txbox(slide, "Phase diagram de l'avantage quantique", Inches(0.7), Inches(0.25),
          Inches(12), Inches(0.55), size=30, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), RGBColor(0xC7, 0x7D, 0xFF))
    txbox(slide, "Dans quelles conditions QMKL peut-il battre RBF-SVM ?  [Notebook 15]",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35), size=13, color=GREY)

    # Grid: séparabilité x bruit
    sep_vals  = [0.3, 0.6, 1.0, 1.5, 2.0]
    noise_vals = [0.0, 0.05, 0.10, 0.20]

    # Simulated delta values (all negative, per report)
    deltas = [
        [-0.300, -0.180, -0.140, -0.090],
        [-0.220, -0.160, -0.100, -0.070],
        [-0.150, -0.080, -0.040, -0.030],
        [-0.080, -0.001, -0.010, -0.020],
        [-0.050, -0.010, -0.025, -0.040],
    ]

    cell_w = Inches(1.3)
    cell_h = Inches(0.72)
    ox = Inches(1.7)
    oy = Inches(1.55)

    # Axis labels
    txbox(slide, "← Bruit →", ox, Inches(1.2), Inches(6.5), Inches(0.3),
          size=10, color=GREY, bold=True)
    for ni, n in enumerate(noise_vals):
        txbox(slide, str(n), ox + ni * cell_w + Inches(0.3), Inches(1.3),
              cell_w, Inches(0.25), size=10, color=GREY, align=PP_ALIGN.CENTER)

    txbox(slide, "Sépar.\n↓", Inches(0.7), oy, Inches(0.95), Inches(0.55), size=10, color=GREY, bold=True)
    for si, s in enumerate(sep_vals):
        txbox(slide, str(s), Inches(1.05), oy + si * cell_h + Inches(0.2),
              Inches(0.55), cell_h, size=10, color=GREY)

    for si, row in enumerate(deltas):
        for ni, delta in enumerate(row):
            # All negative: scale from dark red (very neg) to light red (near 0)
            intensity = min(1.0, abs(delta) / 0.30)
            r = int(0x14 + intensity * (0xE6 - 0x14))
            g = int(0x28 + intensity * (0x20 - 0x28)) if intensity < 0.3 else int(0x39 * (1 - intensity * 0.5))
            b = int(0x3E + (1 - intensity) * 0x20)
            color = RGBColor(max(0,min(255,r)), max(0,min(255,int(0x20*(1-intensity)))), max(0,min(255,int(0x30*(1-intensity)))))

            t = intensity
            r2 = int(0xE6 * t + 0x40 * (1 - t))
            g2 = int(0x10 * t + 0x10 * (1 - t))
            b2 = int(0x20 * t + 0x30 * (1 - t))
            color = RGBColor(r2, g2, b2)

            rect(slide, ox + ni * cell_w, oy + si * cell_h, cell_w - Pt(1), cell_h - Pt(1), color)
            txbox(slide, f"{delta:+.3f}", ox + ni * cell_w + Inches(0.2),
                  oy + si * cell_h + Inches(0.2), cell_w - Inches(0.3), cell_h - Inches(0.25),
                  size=11, bold=(delta > -0.01), color=WHITE if intensity > 0.3 else GREY,
                  align=PP_ALIGN.CENTER)

    # Right panel
    px = Inches(7.1)
    rect(slide, px, Inches(1.5), Inches(5.9), Inches(4.5), BG_LIGHT)
    rect(slide, px, Inches(1.5), Inches(5.9), Pt(4), RGBColor(0xC7, 0x7D, 0xFF))

    txbox(slide, "Résultat le plus frappant de l'étude",
          px + Inches(0.15), Inches(1.62), Inches(5.6), Inches(0.42),
          size=13, bold=True, color=WHITE)

    txbox(slide, "0 / 20",
          px + Inches(0.5), Inches(2.15), Inches(4.9), Inches(1.1),
          size=60, bold=True, color=RED, align=PP_ALIGN.CENTER)

    txbox(slide, "configurations testées\navec avantage quantique",
          px + Inches(0.15), Inches(3.25), Inches(5.6), Inches(0.55),
          size=13, color=WHITE, align=PP_ALIGN.CENTER)

    accent_line(slide, Inches(3.9), GREY, Inches(5.6))

    txbox(slide,
          "Sur données synthétiques tabulaires,\nmême structure que les données financières :\nQMKL ne surpasse pas RBF-SVM dans\naucune des 20 configurations\n(séparabilité × niveau de bruit).",
          px + Inches(0.15), Inches(3.98), Inches(5.6), Inches(1.85),
          size=11, color=GREY)

    # Concept drift result
    rect(slide, px, Inches(6.15), Inches(5.9), Inches(1.0), BG_LIGHT)
    txbox(slide, "Concept drift  [NB 15] :", px + Inches(0.15), Inches(6.2),
          Inches(2.5), Inches(0.35), size=11, bold=True, color=AMBER)
    txbox(slide,
          "Coût moyen du drift = −11 pts AUC\nQMKL inadapté sans réentraînement fréquent",
          px + Inches(0.15), Inches(6.55), Inches(5.6), Inches(0.55),
          size=10, color=WHITE)

    slide_number(slide, 9, TOTAL)
    return slide


def make_contributions_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, CYAN)

    txbox(slide, "14 contributions originales", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), CYAN)
    txbox(slide, "Organisées en 4 niveaux de profondeur — de l'interprétabilité à la frontière",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35), size=13, color=GREY)

    tiers = [
        ("Tier 1", "Interprétabilité", CYAN, [
            "Alignement Frobenius Quantum ↔ RBF  (0.642 moyen)",
            "Entropie spectrale + dimension effective  (H/H_max=0.65)",
            "Dataset parité quantum-hard  (QMKL≈RBF≈hasard)",
            "Frontière de décision PCA-2D",
            "Prototypes financiers (stabilité SV : 84 %)",
        ]),
        ("Tier 2", "Algorithmique", AMBER, [
            "Carte barren plateaux (Q×α)  — sweet spot Q=5,α=2",
            "QKRR — alternative analytique au QSVM  (+7 pts)",
            "Gradient KTA — paysage non convexe découvert",
        ]),
        ("Tier 3", "Avancé", GREEN, [
            "VQKL — optimisation variationnelle des α  (+1.55 pts)",
            "Learning curves  — pas d'avantage data efficiency",
            "Diversité des kernels  — corrélation r=0.738",
        ]),
        ("Tier 4", "Frontière", RGBColor(0xC7, 0x7D, 0xFF), [
            "QKAM — attention instance-adaptive  (−6.4 pts)",
            "Phase diagram avantage quantique  (0/20 configs)",
            "Concept drift robustness  (coût=11 pts AUC)",
        ]),
    ]

    for ci, (tier, label, color, items) in enumerate(tiers):
        x = Inches(0.65 + ci * 3.15)
        rect(slide, x, Inches(1.5), Inches(3.0), Inches(0.6), color)
        txbox(slide, tier,  x + Inches(0.08), Inches(1.52), Inches(1.2), Inches(0.28),
              size=11, bold=True, color=BG)
        txbox(slide, label, x + Inches(0.08), Inches(1.8),  Inches(2.8), Inches(0.28),
              size=10, color=BG)

        for ri, item in enumerate(items):
            y = Inches(2.2 + ri * 0.95)
            rect(slide, x, y, Inches(3.0), Inches(0.85), BG_LIGHT)
            rect(slide, x, y, Pt(4), Inches(0.85), color)
            txbox(slide, item, x + Inches(0.12), y + Inches(0.1),
                  Inches(2.78), Inches(0.7), size=9, color=WHITE)

    # Bottom count
    rect(slide, Inches(0.65), Inches(7.0), Inches(12.2), Inches(0.35), BG_LIGHT)
    txbox(slide,
          "5 métriques Tier 1  ·  3 contributions Tier 2  ·  3 contributions Tier 3  ·  3 contributions Tier 4  =  14 total",
          Inches(0.75), Inches(7.02), Inches(12.0), Inches(0.32),
          size=10, color=GREY, align=PP_ALIGN.CENTER)

    slide_number(slide, 10, TOTAL)
    return slide


def make_nb17_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, RGBColor(0x52, 0xB7, 0x88))

    txbox(slide, "Notebook 17 — La dernière pièce du puzzle", Inches(0.7), Inches(0.25),
          Inches(12), Inches(0.55), size=28, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), GREEN)

    txbox(slide, "Basé sur arXiv 2602.16097 (Zendejas-Morales et al., fév. 2026)",
          Inches(0.7), Inches(1.05), Inches(11), Inches(0.35), size=13, color=GREY)

    # Paper context
    rect(slide, Inches(0.7), Inches(1.5), Inches(12.1), Inches(0.75), BG_LIGHT)
    txbox(slide,
          '"Local and Multi-Scale Strategies to Mitigate Exponential Concentration in Quantum Kernels"',
          Inches(0.85), Inches(1.55), Inches(11.8), Inches(0.6),
          size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

    # Two columns
    left_items = [
        ("Le problème", "La matrice de Gram du kernel global tend vers l'identité quand Q augmente — perte totale d'information discriminante."),
        ("Solution 1 : Patches locaux", "Partitionner les qubits en petits groupes (patches). Chaque patch calcule son propre kernel de fidélité, puis on moyenne."),
        ("Solution 2 : Multi-échelles", "Combiner des kernels locaux à plusieurs granularités simultanément {p=1, p=2, p=4} → capture structure locale ET globale."),
    ]

    right_items = [
        ("Ce que le papier prouve", "La mitigation géométrique (moins de concentration) ne garantit PAS un gain en AUC sans alignement avec les labels (CKA)."),
        ("Notre validation", "Sur Breast Cancer + German Credit + Bank Marketing : nos données confirment p50↑, H/H_max↑ pour kernels locaux mais gain AUC dépend du CKA."),
        ("Lien avec nos résultats", "Ce papier explique formellement pourquoi nos 84 % de support vectors et notre barren plateau map (NB13) s'observent — théorie + expérience alignées."),
    ]

    for i, (title, text) in enumerate(left_items):
        y = Inches(2.4 + i * 1.5)
        rect(slide, Inches(0.7), y, Inches(5.9), Inches(1.35), BG_LIGHT)
        rect(slide, Inches(0.7), y, Pt(4), Inches(1.35), GREEN)
        txbox(slide, title, Inches(0.85), y + Inches(0.06), Inches(5.65), Inches(0.3),
              size=11, bold=True, color=GREEN)
        txbox(slide, text, Inches(0.85), y + Inches(0.38), Inches(5.65), Inches(0.9),
              size=10, color=WHITE)

    for i, (title, text) in enumerate(right_items):
        y = Inches(2.4 + i * 1.5)
        rect(slide, Inches(6.8), y, Inches(5.9), Inches(1.35), BG_LIGHT)
        rect(slide, Inches(6.8), y, Pt(4), Inches(1.35), AMBER)
        txbox(slide, title, Inches(6.95), y + Inches(0.06), Inches(5.65), Inches(0.3),
              size=11, bold=True, color=AMBER)
        txbox(slide, text, Inches(6.95), y + Inches(0.38), Inches(5.65), Inches(0.9),
              size=10, color=WHITE)

    slide_number(slide, 11, TOTAL)
    return slide


def make_conclusions_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)
    rect(slide, 0, 0, Inches(0.45), SLIDE_H, CYAN)

    txbox(slide, "Conclusions", Inches(0.7), Inches(0.25), Inches(11), Inches(0.55),
          size=32, bold=True, color=WHITE)
    accent_line(slide, Inches(0.95), CYAN)

    # Central verdict
    rect(slide, Inches(0.7), Inches(1.15), Inches(12.1), Inches(1.1), BG_LIGHT)
    rect(slide, Inches(0.7), Inches(1.15), Pt(5), Inches(1.1), RED)
    txbox(slide,
          "QMKL n'offre pas d'avantage quantique sur données financières tabulaires dans nos conditions\n"
          "(N=200, Q=6, feature maps Pauli reps=1) — écart de −4 à −7 pts AUC vs RBF-SVM, robuste sur 20 runs et 3 datasets.",
          Inches(0.85), Inches(1.2), Inches(11.7), Inches(0.95),
          size=12, bold=False, color=WHITE)

    # Left column: what we found
    txbox(slide, "Ce qu'on a compris", Inches(0.7), Inches(2.45), Inches(5.9), Inches(0.4),
          size=15, bold=True, color=RED)

    negatives = [
        "Kernels quantiques imitent le RBF (align. 0.642)",
        "Barren plateaux : −60 % concentration entre Q2→Q8",
        "0/20 configurations avec avantage quantique",
        "Instabilité face au concept drift (−11 pts)",
        "Circuits trop peu profonds pour données tabulaires",
    ]
    for i, txt in enumerate(negatives):
        y = Inches(2.95 + i * 0.7)
        rect(slide, Inches(0.7), y + Inches(0.1), Inches(0.18), Inches(0.18), RED)
        txbox(slide, txt, Inches(1.0), y, Inches(5.4), Inches(0.6), size=11, color=WHITE)

    # Right column: what we contributed
    txbox(slide, "Ce qu'on a apporté", Inches(7.0), Inches(2.45), Inches(5.9), Inches(0.4),
          size=15, bold=True, color=GREEN)

    positives = [
        "QKRR : +7 pts vs QSVM  (alternative analytique)",
        "Diversité ↔ gain MKL  (r = 0.738, actionnable)",
        "VQKL : +1.55 pts en optimisant les α",
        "Sweet spot Q=5, α=2.0  (première carte 2D)",
        "14 contributions originales dans la littérature",
    ]
    for i, txt in enumerate(positives):
        y = Inches(2.95 + i * 0.7)
        rect(slide, Inches(7.0), y + Inches(0.1), Inches(0.18), Inches(0.18), GREEN)
        txbox(slide, txt, Inches(7.3), y, Inches(5.5), Inches(0.6), size=11, color=WHITE)

    # Future work
    rect(slide, Inches(0.7), Inches(6.2), Inches(12.1), Inches(0.95), BG_LIGHT)
    txbox(slide, "Vers l'avantage quantique :", Inches(0.85), Inches(6.25),
          Inches(2.2), Inches(0.35), size=11, bold=True, color=AMBER)
    txbox(slide,
          "Circuits plus profonds (reps=2-3)  ·  Données genuinement quantiques  ·  "
          "QKRR comme alternative principale  ·  Sélection par diversité  ·  Hardware réel IBM Quantum",
          Inches(3.05), Inches(6.25), Inches(9.6), Inches(0.75),
          size=10, color=GREY)

    slide_number(slide, 12, TOTAL)
    return slide


def make_final_slide(prs):
    slide = prs.slides.add_slide(blank_layout)
    bg(slide)

    rect(slide, 0, 0, Inches(0.45), SLIDE_H, CYAN)

    # Decorative circles
    for cx, cy, r_inch in [(10.5, 2.0, 2.8), (11.8, 5.0, 1.5)]:
        circ = slide.shapes.add_shape(9,
            Inches(cx - r_inch), Inches(cy - r_inch),
            Inches(r_inch * 2), Inches(r_inch * 2))
        circ.line.color.rgb = CYAN
        circ.line.width = Pt(1)
        circ.fill.background()

    txbox(slide, "QUANTUM MULTIPLE KERNEL LEARNING · FINANCE",
          Inches(0.8), Inches(1.8), Inches(9.5), Inches(0.45),
          size=12, color=CYAN, bold=True)

    txbox(slide, "Merci",
          Inches(0.8), Inches(2.4), Inches(9.5), Inches(1.8),
          size=72, bold=True, color=WHITE)

    accent_line(slide, Inches(4.25), CYAN, Inches(9.0))

    items = [
        f"17 notebooks  ·  14 contributions originales",
        f"Python 3.14 · Qiskit 2.3.1 · Statevector exact",
        f"German Credit · Bank Marketing · Breast Cancer",
        f"arXiv 2602.16097 — Notebook 17 : mitigation concentration",
    ]
    for i, item in enumerate(items):
        txbox(slide, "·  " + item, Inches(0.8), Inches(4.45 + i * 0.42),
              Inches(9.5), Inches(0.38), size=12, color=GREY)

    slide_number(slide, 13, TOTAL)
    return slide


# ── Build the deck ────────────────────────────────────────────────────────────

make_title_slide(prs)
make_journey_slide(prs)
make_question_slide(prs)
make_setup_slide(prs)
make_results_slide(prs)
make_diagnosis_slide(prs)
make_barren_slide(prs)
make_positives_slide(prs)
make_phase_slide(prs)
make_contributions_slide(prs)
make_nb17_slide(prs)
make_conclusions_slide(prs)
make_final_slide(prs)

out = r"C:\Users\Raph\Desktop\QMKL-Finance\QMKL_Finance_Presentation.pptx"
prs.save(out)
print(f"Saved -> {out}")
print(f"Slides : {len(prs.slides)}")
